from flask import Flask, render_template, request, redirect, url_for, flash, jsonify
from flask_uploads import UploadSet, configure_uploads, DOCUMENTS
from pptx import Presentation
import os
import requests
import io
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
import base64
import json

app = Flask(__name__)
app.secret_key = "secret key"

# Configure Flask-Uploads
docs = UploadSet('docs', DOCUMENTS)
app.config['UPLOADED_DOCS_DEST'] = 'uploads'
configure_uploads(app, docs)

# Google Drive setup
SCOPES = ['https://www.googleapis.com/auth/drive.file']
credentials_json = base64.b64decode(os.environ.get('GOOGLE_CREDENTIALS')).decode('utf-8')
creds_dict = json.loads(credentials_json)
creds = service_account.Credentials.from_service_account_info(creds_dict, scopes=SCOPES)
drive_service = build('drive', 'v3', credentials=creds)

def find_fonts(pptx_path):
    presentation = Presentation(pptx_path)
    fonts = set()

    def extract_fonts_from_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name:
                    fonts.add(run.font.name)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                extract_fonts_from_text_frame(shape.text_frame)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            extract_fonts_from_text_frame(cell.text_frame)

            if shape.has_chart:
                chart = shape.chart
                for series in chart.series:
                    for point in series.points:
                        if point.data_label and point.data_label.text_frame:
                            extract_fonts_from_text_frame(point.data_label.text_frame)

    return fonts

def download_font(font_name):
    url = f"https://fonts.google.com/download?family={font_name.replace(' ', '%20')}"
    response = requests.get(url)
    if response.status_code == 200:
        return io.BytesIO(response.content)
    return None

def upload_to_drive(file_path, folder_id=None):
    file_metadata = {'name': os.path.basename(file_path)}
    if folder_id:
        file_metadata['parents'] = [folder_id]
    media = MediaFileUpload(file_path, resumable=True)
    file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    return file.get('id')

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'document' not in request.files:
        flash('No file part')
        return redirect(request.url)
    file = request.files['document']
    if file.filename == '':
        flash('No selected file')
        return redirect(request.url)
    if file:
        file_path = docs.save(file)
        full_path = os.path.join(app.config['UPLOADED_DOCS_DEST'], file_path)
        fonts = find_fonts(full_path)
        font_files = []
        for font in fonts:
            font_file = download_font(font)
            if font_file:
                font_file_path = os.path.join('uploads', f"{font}.zip")
                with open(font_file_path, 'wb') as f:
                    f.write(font_file.read())
                font_files.append(font_file_path)
        # Upload the PPTX file and fonts to Google Drive
        folder_id = request.form.get('folder_id')
        pptx_file_id = upload_to_drive(full_path, folder_id)
        for font_file in font_files:
            upload_to_drive(font_file, folder_id)
        return jsonify({'message': 'File successfully uploaded', 'fonts': list(fonts)})

if __name__ == "__main__":
    app.run(debug=True)

